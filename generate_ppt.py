from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def set_light_background(slide, color=(255, 255, 255)):
    """Set the slide background to a light color (default: white)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def set_text_font(text_frame, font_size):
    """Set the font size for all runs in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

# Create a presentation object
prs = Presentation()

# --- Slide 1: Title Slide ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_light_background(slide)  # Light background
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Emotion-Based Playlist Generator"
subtitle.text = "Using Sentiment Analysis & Spotify API\nYour Name - [Presentation Date]"
set_text_font(title.text_frame, 44)
set_text_font(subtitle.text_frame, 24)

# --- Slide 2: Project Overview ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Project Overview"
content = slide.shapes.placeholders[1]
content.text = (
    "Objective:\nGenerate music playlists based on user input emotions.\n\n"
    "Key Components:\n- Sentiment Analysis using a pre-trained model\n- Spotify API integration\n\n"
    "User Experience:\n- Interactive UI built with Streamlit\n- Scalable backend via FastAPI"
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 3: Concept & Vision ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Concept & Vision"
content = slide.shapes.placeholders[1]
content.text = (
    "Concept:\nCreate an application that analyzes user mood and delivers curated Spotify playlists.\n\n"
    "Vision:\nEnhance the music experience based on emotions using modern NLP and API integrations."
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 4: Tech Stack & Tools ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Tech Stack & Tools"
content = slide.shapes.placeholders[1]
content.text = (
    "Libraries:\n• Hugging Face Transformers\n• Spotipy (Spotify API)\n• NLTK\n• scikit-learn\n\n"
    "Frameworks:\n• Streamlit (UI)\n• FastAPI (Backend)\n\n"
    "Dataset:\n• spotify_millsongdata.csv\n\n"
    "Additional Tools:\n• python-dotenv for configuration"
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 5: Repository Structure ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Repository Structure"
content = slide.shapes.placeholders[1]
content.text = (
    "Emotion_Based_Playlist_Generator/\n"
    "├── data/\n"
    "│   └── spotify_millsongdata.csv\n"
    "├── src/\n"
    "│   ├── __init__.py\n"
    "│   ├── sentiment_analysis.py\n"
    "│   ├── spotify_recommender.py\n"
    "│   └── utils.py\n"
    "├── app.py\n"
    "├── main.py\n"
    "├── config.py\n"
    "├── requirements.txt\n"
    "└── README.md"
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 16)

# --- Slide 6: Sentiment Analysis Module ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Sentiment Analysis Module"
content = slide.shapes.placeholders[1]
content.text = (
    "File: src/sentiment_analysis.py\n\n"
    "• Uses a pre-trained Hugging Face model (e.g., cardiffnlp/twitter-roberta-base-sentiment).\n"
    "• Maps output labels (LABEL_0 → NEGATIVE, LABEL_1 → NEUTRAL, LABEL_2 → POSITIVE).\n"
    "• Returns sentiment label and confidence score."
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 7: Spotify Integration Module ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Spotify Integration Module"
content = slide.shapes.placeholders[1]
content.text = (
    "File: src/spotify_recommender.py\n\n"
    "• Uses Spotipy with Spotify API credentials from config.py.\n"
    "• Maps sentiment to query keywords (e.g., 'feel good' for POSITIVE, 'sad songs' for NEGATIVE).\n"
    "• Retrieves a valid playlist URL based on the detected mood."
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 8: User Interface (Streamlit) ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "User Interface (Streamlit)"
content = slide.shapes.placeholders[1]
content.text = (
    "File: app.py\n\n"
    "• Users enter mood or thoughts.\n"
    "• Displays detected sentiment with confidence.\n"
    "• Provides a clickable link to the recommended Spotify playlist.\n"
    "• Clean and interactive design."
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 9: Backend API (FastAPI) ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Backend API (FastAPI)"
content = slide.shapes.placeholders[1]
content.text = (
    "File: main.py\n\n"
    "• Endpoints:\n   - /analyze/ returns sentiment analysis for input text.\n"
    "   - /playlist/{mood} returns a Spotify playlist URL based on mood.\n\n"
    "• Supports integration with other apps and services."
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 10: Challenges & Future Work ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Challenges & Future Work"
content = slide.shapes.placeholders[1]
content.text = (
    "Challenges:\n• Misclassification of ambiguous text\n• Query mapping inconsistencies\n\n"
    "Future Improvements:\n• Explore alternative sentiment models\n• Fine-tune query keywords\n• Enhance UI/UX based on feedback"
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 11: Demo & Conclusion ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Demo & Conclusion"
content = slide.shapes.placeholders[1]
content.text = (
    "Demo:\n• Live demonstration of the Streamlit app and API endpoints\n\n"
    "Conclusion:\n• Blends modern NLP with real-time music recommendations\n• Potential for a full-scale emotion-driven music platform"
)
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# --- Slide 12: Q & A ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_background(slide)
slide.shapes.title.text = "Q & A"
content = slide.shapes.placeholders[1]
content.text = "Thank You!\nQuestions?"
set_text_font(slide.shapes.title.text_frame, 32)
set_text_font(content.text_frame, 18)

# Save the presentation
pptx_filename = "Emotion_Based_Playlist_Generator.pptx"
prs.save(pptx_filename)
print(f"Presentation saved as '{pptx_filename}'")
